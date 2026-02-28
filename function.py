import io
import os
import re
import shutil
import subprocess
import tempfile
import time
from datetime import datetime

from pptx import Presentation

from dotenv import load_dotenv
from openai import OpenAI
from pydub import AudioSegment
from pydub.silence import split_on_silence
import streamlit as st
from scipy.signal import butter, lfilter
import numpy as np

from constants import (
    MINUTES_FILE_NAME,
    OPENAI_SUMMARY_MODEL,
    OPENAI_SYSTEM_PROMPT,
    OPENAI_WHISPER_MODEL,
    RAW_TRANSCRIPTION_FILE_NAME,
    SESSION_DEFAULTS,
    SUMMARY_PROMPT_TEMPLATE,
    TRANSCRIPTION_WORD_REPLACEMENTS,
    WHISPER_LANGUAGE,
    WHISPER_MAX_FILE_MB,
    WHISPER_PROMPT,
)

load_dotenv()


def initialize_session_state():
    for key, value in SESSION_DEFAULTS.items():
        if key not in st.session_state:
            st.session_state[key] = value


def clear_session_state():
    st.session_state.summary = None
    st.session_state.transcription = None
    st.session_state.save_success = False
    st.session_state.saved_file_path = None
    st.session_state.save_error = None
    st.session_state.transcription_save_success = False
    st.session_state.saved_transcription_file_path = None
    st.session_state.transcription_save_error = None
    st.session_state.meeting_info_input = ""
    st.session_state.uploader_version += 1


def get_openai_client():
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        return None
    return OpenAI(api_key=api_key)


def extract_audio_from_video(video_path, output_audio_path):
    try:
        ffmpeg_bin = shutil.which("ffmpeg")
        if not ffmpeg_bin:
            return False, "ffmpeg が見つかりません。PATH を確認してください。"

        ffprobe_bin = shutil.which("ffprobe")
        if ffprobe_bin:
            probe_command = [
                ffprobe_bin,
                "-v",
                "error",
                "-select_streams",
                "a",
                "-show_entries",
                "stream=index",
                "-of",
                "csv=p=0",
                video_path,
            ]
            probe_result = subprocess.run(probe_command, capture_output=True, text=True)
            if probe_result.returncode == 0 and not probe_result.stdout.strip():
                return False, "この動画ファイルには音声トラックがありません。音声付きの mp4 または mp3/m4a をアップロードしてください。"

        command = [
            ffmpeg_bin,
            "-y",
            "-i",
            video_path,
            "-vn",
            "-acodec",
            "libmp3lame",
            output_audio_path,
        ]

        result = subprocess.run(command, capture_output=True, text=True)
        if result.returncode != 0:
            error_message = (result.stderr or result.stdout or "ffmpeg command failed").strip()
            if "does not contain any stream" in error_message.lower():
                return False, "この動画ファイルには音声トラックがありません。音声付きの mp4 または mp3/m4a をアップロードしてください。"
            return False, error_message

        print(f"Audio extracted and saved to {output_audio_path}")
        return True, None
    except Exception as e:
        print(f"Error extracting audio: {e}")
        return False, str(e)


def split_audio(audio_path, output_dir, silence_thresh=-40, min_silence_len=700):
    """
    Splits an audio file into chunks based on silence, with noise reduction applied.

    Args:
        audio_path (str): Path to the input audio file.
        output_dir (str): Directory to save the audio chunks.
        silence_thresh (int): Silence threshold in dBFS.
        min_silence_len (int): Minimum length of silence to consider for splitting (in ms).
    """
    try:
        audio = AudioSegment.from_file(audio_path)

        # Apply noise reduction
        noise_reduced_audio = reduce_noise(audio)

        chunks = split_on_silence(
            noise_reduced_audio,
            min_silence_len=min_silence_len,
            silence_thresh=silence_thresh,
        )

        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        for i, chunk in enumerate(chunks):
            chunk_path = os.path.join(output_dir, f"chunk_{i}.mp3")
            chunk.export(chunk_path, format="mp3")
            print(f"Exported chunk {i} to {chunk_path}")
    except Exception as e:
        print(f"Error splitting audio: {e}")


def transcribe_audio_with_whisper(audio_path):
    try:
        client = get_openai_client()
        if not client:
            msg = "OPENAI_API_KEY is not set"
            print(f"Error during transcription: {msg}")
            return None, msg

        file_size_bytes = os.path.getsize(audio_path)
        if file_size_bytes <= WHISPER_MAX_FILE_MB * 1024 * 1024:
            text, single_error = transcribe_single_file(client, audio_path)
            return text, single_error

        chunk_dir = tempfile.mkdtemp(prefix="whisper_chunks_")
        try:
            chunk_paths, split_error = split_audio_for_whisper_limit(
                audio_path,
                chunk_dir,
                WHISPER_MAX_FILE_MB,
            )
            if split_error:
                print(f"Error during chunk split: {split_error}")
                return None, f"音声分割に失敗しました: {split_error}"

            transcripts = []
            for chunk_path in chunk_paths:
                chunk_text = transcribe_single_file(client, chunk_path)
                chunk_text, chunk_error = transcribe_single_file(client, chunk_path)
                if not chunk_text:
                    print(f"Error during transcription for chunk: {chunk_path}, {chunk_error}")
                    return None, f"分割音声の文字起こしに失敗しました: {chunk_error}"
                transcripts.append(chunk_text)

            return "\n".join(transcripts), None
        finally:
            shutil.rmtree(chunk_dir, ignore_errors=True)
    except Exception as e:
        print(f"Error during transcription: {e}")
        return None, str(e)


def apply_word_replacements(text):
    """
    Applies word replacements defined in TRANSCRIPTION_WORD_REPLACEMENTS
    to the transcribed text.

    Args:
        text (str): The transcribed text.

    Returns:
        str: The text with replacements applied.
    """
    for wrong, correct in TRANSCRIPTION_WORD_REPLACEMENTS.items():
        text = text.replace(wrong, correct)
    return text


def transcribe_single_file(client, audio_path):
    try:
        with open(audio_path, "rb") as audio_file:
            response = client.audio.transcriptions.create(
                model=OPENAI_WHISPER_MODEL,
                file=audio_file,
                language=WHISPER_LANGUAGE,
                prompt=WHISPER_PROMPT,
            )
        text = apply_word_replacements(response.text)
        return text, None
    except Exception as e:
        return None, str(e)


def split_audio_for_whisper_limit(audio_path, output_dir, max_chunk_mb):
    try:
        audio = AudioSegment.from_file(audio_path)
        if len(audio) == 0:
            return [], "音声が空のため分割できません。"

        max_chunk_bytes = max_chunk_mb * 1024 * 1024
        source_size = os.path.getsize(audio_path)
        bytes_per_ms = source_size / max(len(audio), 1)
        estimated_chunk_ms = int((max_chunk_bytes * 0.9) / max(bytes_per_ms, 1e-6))
        chunk_ms = max(estimated_chunk_ms, 30_000)

        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        chunk_paths = []
        start_ms = 0
        chunk_index = 0

        while start_ms < len(audio):
            end_ms = min(start_ms + chunk_ms, len(audio))
            segment = audio[start_ms:end_ms]
            chunk_path = os.path.join(output_dir, f"chunk_{chunk_index}.mp3")
            segment.export(chunk_path, format="mp3", bitrate="128k")

            while os.path.getsize(chunk_path) > max_chunk_bytes and (end_ms - start_ms) > 15_000:
                end_ms = start_ms + int((end_ms - start_ms) * 0.8)
                segment = audio[start_ms:end_ms]
                segment.export(chunk_path, format="mp3", bitrate="128k")

            if os.path.getsize(chunk_path) > max_chunk_bytes:
                return [], f"チャンク作成に失敗しました: {chunk_index}"

            chunk_paths.append(chunk_path)
            start_ms = end_ms
            chunk_index += 1

        return chunk_paths, None
    except Exception as e:
        return [], str(e)


def extract_text_from_pptx(file_bytes):
    """
    PowerPointファイルのバイト列からテキストを抽出します。

    Args:
        file_bytes (bytes): PowerPointファイルのバイト列。

    Returns:
        tuple[str | None, str | None]: (抽出テキスト, エラーメッセージ)
    """
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            slide_texts.append(line)
            if slide_texts:
                texts.append(f"【スライド{slide_num}】\n" + "\n".join(slide_texts))
        return "\n\n".join(texts), None
    except Exception as e:
        return None, str(e)


def extract_text_from_pptx_files(uploaded_files):
    """
    複数のPowerPointファイルからテキストを結合して抽出します。

    Args:
        uploaded_files (list): StreamlitのUploadedFileオブジェクトのリスト（最大7件）。

    Returns:
        tuple[str, list[str]]: (結合テキスト, エラーメッセージのリスト)
    """
    all_texts = []
    errors = []
    for i, f in enumerate(uploaded_files, start=1):
        text, error = extract_text_from_pptx(f.getvalue())
        if error:
            errors.append(f"{f.name}: {error}")
        elif text:
            all_texts.append(f"=== 資料{i}: {f.name} ===\n{text}")
    return "\n\n".join(all_texts), errors


def summarize_transcription(transcription, meeting_info="", reference_text=""):
    try:
        client = get_openai_client()
        if not client:
            print("Error during summarization: OPENAI_API_KEY is not set")
            return None

        prompt = SUMMARY_PROMPT_TEMPLATE.format(
            meeting_info=meeting_info or "（入力なし）",
            reference_material=reference_text or "（アップロードなし）",
            transcription=transcription,
        )

        response = client.chat.completions.create(
            model=OPENAI_SUMMARY_MODEL,
            messages=[
                {"role": "system", "content": OPENAI_SYSTEM_PROMPT},
                {"role": "user", "content": prompt},
            ],
        )

        summary_text = response.choices[0].message.content
        if not summary_text:
            return None

        fallback_meeting = parse_meeting_basic_info(meeting_info)
        summary = {
            "meeting_name": fallback_meeting["meeting_name"],
            "meeting_datetime": fallback_meeting["meeting_datetime"],
            "participants": fallback_meeting["participants"],
            "location_url": fallback_meeting["location_url"],
            "agenda": "",
            "main_points": "",
            "decisions": "",
        }

        if "0. 会議基本情報" in summary_text and "1. 議題の説明:" in summary_text:
            meeting_info_text = summary_text.split("0. 会議基本情報")[1].split("1. 議題の説明:")[0].strip()
            parsed_meeting = parse_meeting_basic_info(meeting_info_text)
            summary["meeting_name"] = parsed_meeting["meeting_name"]
            summary["meeting_datetime"] = parsed_meeting["meeting_datetime"]
            summary["participants"] = parsed_meeting["participants"]
            summary["location_url"] = parsed_meeting["location_url"]

        if "1. 議題の説明:" in summary_text:
            summary["agenda"] = summary_text.split("1. 議題の説明:")[1].split("2. 主な発言:")[0].strip()
        if "2. 主な発言:" in summary_text:
            summary["main_points"] = summary_text.split("2. 主な発言:")[1].split("3. 決定事項:")[0].strip()
        if "3. 決定事項:" in summary_text:
            summary["decisions"] = summary_text.split("3. 決定事項:")[1].strip()

        return summary
    except Exception as e:
        print(f"Error during summarization: {e}")
        return None


def parse_meeting_basic_info(text):
    default_info = {
        "meeting_name": "不明",
        "meeting_datetime": "不明",
        "participants": "不明",
        "location_url": "不明",
    }

    if not text:
        return default_info

    patterns = {
        "meeting_name": r"会議名\s*[:：]\s*(.+)",
        "meeting_datetime": r"日時\s*[:：]\s*(.+)",
        "participants": r"参加者\s*[:：]\s*(.+)",
        "location_url": r"場所\s*/\s*URL\s*[:：]\s*(.+)",
    }

    parsed = default_info.copy()
    for key, pattern in patterns.items():
        match = re.search(pattern, text)
        if match and match.group(1).strip():
            parsed[key] = match.group(1).strip()

    return parsed


def export_to_local_folder(summary, output_dir):
    try:
        output_dir_abs = os.path.abspath(output_dir)
        if not os.path.exists(output_dir_abs):
            os.makedirs(output_dir_abs)

        file_path = os.path.join(output_dir_abs, MINUTES_FILE_NAME)
        with open(file_path, "w", encoding="utf-8") as file:
            file.write(build_minutes_text(summary))

        return file_path, None
    except Exception as e:
        print(f"Error exporting to local folder: {e}")
        return None, str(e)


def export_transcription_to_local_folder(transcription, output_dir):
    try:
        output_dir_abs = os.path.abspath(output_dir)
        if not os.path.exists(output_dir_abs):
            os.makedirs(output_dir_abs)

        file_path = os.path.join(output_dir_abs, RAW_TRANSCRIPTION_FILE_NAME)
        with open(file_path, "w", encoding="utf-8") as file:
            file.write(build_transcription_raw_text(transcription))

        return file_path, None
    except Exception as e:
        print(f"Error exporting transcription to local folder: {e}")
        return None, str(e)


def build_transcription_raw_text(transcription):
    generated_at = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    return (
        f"生成日時: {generated_at}\n"
        "\n"
        "--- 文字起こし生データ ---\n"
        f"{transcription or ''}\n"
    )


def build_minutes_text(summary):
    return (
        "会議基本情報:\n"
        f"会議名: {summary.get('meeting_name', '不明')}\n"
        f"日時: {summary.get('meeting_datetime', '不明')}\n"
        f"参加者: {summary.get('participants', '不明')}\n"
        f"場所/URL: {summary.get('location_url', '不明')}\n\n"
        f"議題の説明:\n{summary['agenda']}\n\n"
        f"主な発言:\n{summary['main_points']}\n\n"
        f"決定事項:\n{summary['decisions']}\n"
    )


def cleanup_old_files(directory, retention_period_days=90):
    try:
        current_time = time.time()
        retention_period_seconds = retention_period_days * 24 * 66 * 60 * 60

        if not os.path.exists(directory):
            print(f"Directory {directory} does not exist.")
            return

        for filename in os.listdir(directory):
            file_path = os.path.join(directory, filename)
            if os.path.isfile(file_path):
                file_age = current_time - os.path.getmtime(file_path)
                if file_age > retention_period_seconds:
                    os.remove(file_path)
                    print(f"Deleted old file: {file_path}")
    except Exception as e:
        print(f"Error during cleanup: {e}")


def process_vtt_file(vtt_path):
    try:
        with open(vtt_path, "r", encoding="utf-8") as file:
            lines = file.readlines()

        text_lines = []
        for line in lines:
            line = line.strip()
            if not line or "-->" in line or line.isdigit():
                continue
            text_lines.append(line)

        return "\n".join(text_lines)
    except Exception as e:
        print(f"Error processing .vtt file: {e}")
        return None


def reduce_noise(audio_segment, noise_reduction_level=0.02):
    """
    Reduces noise from an audio segment using a simple low-pass filter.

    Args:
        audio_segment (AudioSegment): The input audio segment.
        noise_reduction_level (float): The level of noise reduction (0.0 to 1.0).

    Returns:
        AudioSegment: The noise-reduced audio segment.
    """
    # Convert AudioSegment to numpy array
    samples = np.array(audio_segment.get_array_of_samples())

    # Design a low-pass filter
    def butter_lowpass(cutoff, fs, order=5):
        nyquist = 0.5 * fs
        normal_cutoff = cutoff / nyquist
        b, a = butter(order, normal_cutoff, btype="low", analog=False)
        return b, a

    def apply_lowpass_filter(data, cutoff_freq, sample_rate):
        b, a = butter_lowpass(cutoff_freq, sample_rate)
        return lfilter(b, a, data)

    # Apply low-pass filter
    sample_rate = audio_segment.frame_rate
    cutoff_frequency = sample_rate * noise_reduction_level
    filtered_samples = apply_lowpass_filter(samples, cutoff_frequency, sample_rate)

    # Convert back to AudioSegment
    reduced_audio = audio_segment._spawn(filtered_samples.astype(np.int16).tobytes())
    return reduced_audio

# Example usage in split_audio or other functions:
# audio = AudioSegment.from_file(audio_path)
# noise_reduced_audio = reduce_noise(audio)
# chunks = split_on_silence(noise_reduced_audio, ...)
